"""
解析几何题目生成系统 - PPT演示文稿（Claude Design System）
Analytic Geometry Problem Generator - PPT Presentation (Anthropic Design System)

Color palette & typography from DESIGN.md:
- Canvas: #faf9f5 (warm cream)
- Primary/CTA: #cc785c (coral)
- Surface Card: #efe9de (cream card)
- Surface Dark: #181715 (dark navy)
- Ink: #141413 (warm black)
- Muted: #6c6a64
- Accent Teal: #5db8a6
- Accent Amber: #e8a55a

Typography: serif display (Tiemos/Garamond fallback) + sans body (Inter/system)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os


class GeometryPPTCreator:
    """PPT演示文稿创建器 — Claude Design System"""

    # ==================== DESIGN TOKENS ====================
    # From DESIGN.md

    # Colors
    CANVAS = RGBColor(0xFA, 0xF9, 0xF5)           # #faf9f5  warm cream canvas
    SURFACE_CARD = RGBColor(0xEF, 0xE9, 0xDE)      # #efe9de  cream card
    SURFACE_SOFT = RGBColor(0xF5, 0xF0, 0xE8)      # #f5f0e8  soft cream
    SURFACE_DARK = RGBColor(0x18, 0x17, 0x15)      # #181715  dark navy
    SURFACE_DARK_ELEVATED = RGBColor(0x25, 0x23, 0x20)  # #252320
    SURFACE_DARK_SOFT = RGBColor(0x1F, 0x1E, 0x1B)  # #1f1e1b

    PRIMARY = RGBColor(0xCC, 0x78, 0x5C)            # #cc785c  coral
    PRIMARY_ACTIVE = RGBColor(0xA9, 0x58, 0x3E)     # #a9583e
    ON_PRIMARY = RGBColor(0xFF, 0xFF, 0xFF)          # white

    INK = RGBColor(0x14, 0x14, 0x13)                # #141413  warm black
    BODY = RGBColor(0x3D, 0x3D, 0x3A)               # #3d3d3a
    BODY_STRONG = RGBColor(0x25, 0x25, 0x23)        # #252523
    MUTED = RGBColor(0x6C, 0x6A, 0x64)              # #6c6a64
    MUTED_SOFT = RGBColor(0x8E, 0x8B, 0x82)         # #8e8b82
    HAIRLINE = RGBColor(0xE6, 0xDF, 0xD8)           # #e6dfd8

    ON_DARK = RGBColor(0xFA, 0xF9, 0xF5)            # cream on dark
    ON_DARK_SOFT = RGBColor(0xA0, 0x9D, 0x96)       # #a09d96

    ACCENT_TEAL = RGBColor(0x5D, 0xB8, 0xA6)        # #5db8a6
    ACCENT_AMBER = RGBColor(0xE8, 0xA5, 0x5A)       # #e8a55a
    SUCCESS = RGBColor(0x5D, 0xB8, 0x72)            # #5db872

    # Typography tokens — fallback fonts for PPT
    # Display: Georgia (serif) for Copernicus/Tiempos Headline
    # Body: Calibri (humanist sans) for StyreneB/Inter
    DISPLAY_FONT = "Georgia"
    BODY_FONT = "Calibri"
    CODE_FONT = "Consolas"

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9
        self.prs.slide_height = Inches(7.5)

    # ==================== SLIDE BUILDERS ====================

    def _set_background(self, slide, color):
        """Set solid background color"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_text(self, slide, left, top, width, height,
                  text, font_name=None, font_size=Pt(16),
                  font_color=None, bold=False, italic=False,
                  alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                  line_spacing=1.55, space_after=Pt(8)):
        """Add a text box with styling"""
        if font_name is None:
            font_name = self.BODY_FONT
        if font_color is None:
            font_color = self.BODY

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text
        p.alignment = alignment
        p.line_spacing = line_spacing
        p.space_after = space_after

        run = p.runs[0]
        run.font.name = font_name
        run.font.size = font_size
        run.font.color.rgb = font_color
        run.font.bold = bold
        run.font.italic = italic

        return textbox, tf

    def _add_bullet_list(self, slide, left, top, width, height,
                         items, font_name=None, font_size=Pt(16),
                         font_color=None, line_spacing=1.5,
                         bullet_char="•", space_after=Pt(6)):
        """Add a bulleted list"""
        if font_name is None:
            font_name = self.BODY_FONT
        if font_color is None:
            font_color = self.BODY

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = f"{bullet_char}  {item}"
            p.line_spacing = line_spacing
            p.space_after = space_after

            run = p.runs[0]
            run.font.name = font_name
            run.font.size = font_size
            run.font.color.rgb = font_color

        return textbox

    def _add_heading(self, slide, left, top, width, height,
                     text, size="lg", color=None):
        """Add a display heading (serif)"""
        size_map = {
            "xl": (Pt(64), -1.5),
            "lg": (Pt(48), -1.0),
            "md": (Pt(36), -0.5),
            "sm": (Pt(28), -0.3),
        }
        font_size, tracking = size_map.get(size, (Pt(48), -1.0))
        if color is None:
            color = self.INK

        return self._add_text(
            slide, left, top, width, height,
            text, font_name=self.DISPLAY_FONT,
            font_size=font_size, font_color=color,
            bold=False, line_spacing=1.1, space_after=Pt(4)
        )

    def _add_subheading(self, slide, left, top, width, height,
                        text, size="md", color=None):
        """Add a sans subheading"""
        size_map = {
            "lg": Pt(22),
            "md": Pt(18),
            "sm": Pt(16),
        }
        font_size = size_map.get(size, Pt(18))
        if color is None:
            color = self.BODY_STRONG

        return self._add_text(
            slide, left, top, width, height,
            text, font_name=self.BODY_FONT,
            font_size=font_size, font_color=color,
            bold=True, line_spacing=1.4, space_after=Pt(6)
        )

    def _add_coral_band(self, slide, top=0, height=Inches(1.2)):
        """Add a coral top band"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), top,
            Inches(13.333), height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.PRIMARY
        shape.line.fill.background()
        return shape

    def _add_dark_band(self, slide, top=0, height=Inches(1.2)):
        """Add a dark navy top band"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), top,
            Inches(13.333), height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.SURFACE_DARK
        shape.line.fill.background()
        return shape

    def _add_card(self, slide, left, top, width, height,
                  bg_color=None, border=True):
        """Add a cream card with optional hairline border"""
        if bg_color is None:
            bg_color = self.SURFACE_CARD

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border:
            shape.line.color.rgb = self.HAIRLINE
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()

        return shape

    def _add_dark_card(self, slide, left, top, width, height):
        """Add a dark navy card (product mockup)"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.SURFACE_DARK
        shape.line.fill.background()
        return shape

    # ==================== SLIDE TYPES ====================

    def _slide_title(self):
        """Cover slide: cream canvas, coral wordmark accent"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide, self.CANVAS)

        # Coral accent stripe (thin, left-aligned)
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(2.2),
            Inches(0.08), Inches(2.8)
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = self.PRIMARY
        stripe.line.fill.background()

        # Title (serif display, large, negative tracking simulated)
        self._add_heading(slide, Inches(1.4), Inches(2.2), Inches(10), Inches(1.8),
                          "解析几何题目生成系统",
                          size="xl", color=self.INK)

        # Subtitle
        self._add_text(slide, Inches(1.4), Inches(4.2), Inches(9), Inches(0.8),
                       "基于数学参数的动态题目生成与精确配图渲染",
                       font_size=Pt(22), font_color=self.MUTED,
                       line_spacing=1.4)

        # Coral CTA badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.4), Inches(5.4),
            Inches(3.2), Inches(0.55)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = self.PRIMARY
        badge.line.fill.background()
        self._add_text(slide, Inches(1.4), Inches(5.42), Inches(3.2), Inches(0.5),
                       "Anthropic Design System",
                       font_size=Pt(13), font_color=self.ON_PRIMARY,
                       bold=True, alignment=PP_ALIGN.CENTER)

        return slide

    def _slide_toc(self):
        """Table of contents: cream card list"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide, self.CANVAS)

        self._add_heading(slide, Inches(1), Inches(0.8), Inches(11), Inches(1),
                          "目录", size="lg", color=self.INK)

        # Coral divider
        div = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(1.8), Inches(2), Inches(0.04)
        )
        div.fill.solid()
        div.fill.fore_color.rgb = self.PRIMARY
        div.line.fill.background()

        sections = [
            ("01", "系统概述与设计理念"),
            ("02", "系统架构"),
            ("03", "核心算法：题目生成引擎"),
            ("04", "核心算法：精确配图渲染"),
            ("05", "演示：椭圆题目"),
            ("06", "演示：双曲线题目"),
            ("07", "演示：抛物线题目"),
            ("08", "演示：极坐标题目"),
            ("09", "技术亮点与扩展方向"),
        ]

        for i, (num, title) in enumerate(sections):
            y = Inches(2.2) + Inches(i * 0.55)
            # Number badge
            self._add_text(slide, Inches(1), y, Inches(0.8), Inches(0.45),
                           num, font_name=self.DISPLAY_FONT,
                           font_size=Pt(20), font_color=self.PRIMARY,
                           bold=True, line_spacing=1.0)
            # Section title
            self._add_text(slide, Inches(1.9), y, Inches(9), Inches(0.45),
                           title, font_size=Pt(18), font_color=self.BODY_STRONG,
                           line_spacing=1.0)

        return slide

    def _slide_content(self, title, bullets_left=None, bullets_right=None,
                       image_path=None, band_color="coral"):
        """Generic content slide with cream canvas + optional cards"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide, self.CANVAS)

        # Top band
        if band_color == "coral":
            self._add_coral_band(slide)
        else:
            self._add_dark_band(slide)

        # Band title
        title_color = self.ON_DARK if band_color == "dark" else self.ON_PRIMARY
        self._add_text(slide, Inches(0.8), Inches(0.2), Inches(11.5), Inches(0.8),
                       title, font_name=self.DISPLAY_FONT,
                       font_size=Pt(36), font_color=title_color,
                       bold=False, line_spacing=1.15, space_after=Pt(0))

        if image_path and os.path.exists(image_path):
            # Image + text layout
            if bullets_right is None:
                bullets_right = []

            # Image in dark card
            card = self._add_dark_card(slide, Inches(7.2), Inches(1.6), Inches(5.5), Inches(5.3))
            slide.shapes.add_picture(
                image_path,
                Inches(7.4), Inches(1.8), Inches(5.1), Inches(4.9)
            )

            # Text in cream card
            if bullets_left:
                text_card = self._add_card(slide, Inches(0.6), Inches(1.6), Inches(6.2), Inches(5.3))
                self._add_bullet_list(slide, Inches(0.9), Inches(1.9), Inches(5.6), Inches(4.8),
                                      bullets_left, font_size=Pt(15), line_spacing=1.55)
        elif bullets_left and bullets_right:
            # Two-column cards
            left_card = self._add_card(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.3))
            self._add_bullet_list(slide, Inches(0.9), Inches(1.9), Inches(5.2), Inches(4.8),
                                  bullets_left, font_size=Pt(15), line_spacing=1.5)

            right_card = self._add_card(slide, Inches(6.9), Inches(1.6), Inches(5.8), Inches(5.3))
            self._add_bullet_list(slide, Inches(7.2), Inches(1.9), Inches(5.2), Inches(4.8),
                                  bullets_right, font_size=Pt(15), line_spacing=1.5)
        elif bullets_left:
            # Single wide card
            text_card = self._add_card(slide, Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.3))
            self._add_bullet_list(slide, Inches(0.9), Inches(1.9), Inches(11.5), Inches(4.8),
                                  bullets_left, font_size=Pt(15), line_spacing=1.55)

        return slide

    def _slide_images(self, title, images, captions=None):
        """Image showcase slide with dark cards"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide, self.CANVAS)

        # Dark top band
        self._add_dark_band(slide)
        self._add_text(slide, Inches(0.8), Inches(0.2), Inches(11.5), Inches(0.8),
                       title, font_name=self.DISPLAY_FONT,
                       font_size=Pt(36), font_color=self.ON_DARK,
                       bold=False, line_spacing=1.15)

        # Layout: up to 3 images in dark cards
        n = len(images)
        if n == 1:
            positions = [(Inches(2.5), Inches(1.6), Inches(8.3), Inches(5.3))]
        elif n == 2:
            positions = [
                (Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.3)),
                (Inches(6.9), Inches(1.6), Inches(5.8), Inches(5.3)),
            ]
        else:
            positions = [
                (Inches(0.6), Inches(1.6), Inches(3.8), Inches(5.3)),
                (Inches(4.7), Inches(1.6), Inches(3.8), Inches(5.3)),
                (Inches(8.9), Inches(1.6), Inches(3.8), Inches(5.3)),
            ]

        for i, img_path in enumerate(images[:3]):
            if i < len(positions) and os.path.exists(img_path):
                left, top, w, h = positions[i]
                card = self._add_dark_card(slide, left, top, w, h)
                # Image inside card with padding
                slide.shapes.add_picture(
                    img_path,
                    left + Inches(0.15), top + Inches(0.15),
                    w - Inches(0.3), h - Inches(0.8)
                )
                # Caption
                if captions and i < len(captions):
                    self._add_text(slide, left, top + h - Inches(0.55), w, Inches(0.45),
                                   captions[i], font_size=Pt(12),
                                   font_color=self.MUTED_SOFT,
                                   alignment=PP_ALIGN.CENTER, line_spacing=1.0)

        return slide

    def _slide_cta_coral(self):
        """Coral CTA card slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide, self.CANVAS)

        # Full coral card
        card = self._add_card(slide, Inches(0.8), Inches(1.2),
                              Inches(11.7), Inches(5.1),
                              bg_color=self.PRIMARY, border=False)

        self._add_text(slide, Inches(1.6), Inches(2.0), Inches(10), Inches(1.5),
                       "谢谢观看",
                       font_name=self.DISPLAY_FONT,
                       font_size=Pt(64), font_color=self.ON_PRIMARY,
                       bold=False, line_spacing=1.05,
                       alignment=PP_ALIGN.CENTER)

        self._add_text(slide, Inches(1.6), Inches(3.8), Inches(10), Inches(1),
                       "解析几何题目生成系统 · 让数学出题更高效",
                       font_size=Pt(22), font_color=RGBColor(0xFF, 0xE0, 0xD0),
                       alignment=PP_ALIGN.CENTER, line_spacing=1.3)

        # Cream button on coral
        btn = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.2), Inches(5.0), Inches(2.9), Inches(0.5)
        )
        btn.fill.solid()
        btn.fill.fore_color.rgb = self.CANVAS
        btn.line.fill.background()
        self._add_text(slide, Inches(5.2), Inches(5.02), Inches(2.9), Inches(0.45),
                       "Generated with ❤",
                       font_size=Pt(14), font_color=self.PRIMARY,
                       bold=True, alignment=PP_ALIGN.CENTER)

        return slide

    # ==================== BUILD PRESENTATION ====================

    def create_presentation(self, output_path: str = None) -> str:
        """Build the full presentation"""
        if output_path is None:
            output_path = "/root/analytic_geometry_generator/output/解析几何题目生成系统.pptx"

        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        # ---- Slide 1: Title ----
        self._slide_title()

        # ---- Slide 2: TOC ----
        self._slide_toc()

        # ---- Slide 3: Overview ----
        self._slide_content("1. 系统概述与设计理念", bullets_left=[
            "核心目标",
            "  根据用户指定的知识点与难度等级，动态生成解析几何题目",
            "  题目必须能够被证明或被求解（数学可解性保证）",
            "  配图中所有关键点、直线、曲线位置与题干数学参数严格一致",
            "",
            "支持的知识点",
            "  椭圆 (Ellipse)：标准方程、焦点弦、焦点三角形",
            "  双曲线 (Hyperbola)：渐近线、焦点弦、焦点三角形",
            "  抛物线 (Parabola)：焦点弦、焦点弦性质证明",
            "  极坐标 (Polar)：坐标互化、圆锥曲线极坐标方程",
        ], bullets_right=[
            "难度等级",
            "  ★☆☆ Level 1 基础：标准方程、焦点、顶点",
            "  ★★☆ Level 2 进阶：弦长、面积、位置关系",
            "  ★★★ Level 3 竞赛：综合证明、最值问题",
            "",
            "输出格式",
            "  LaTeX 格式题干与完整解答",
            "  Matplotlib 精确几何配图 (PNG)",
            "  PowerPoint 演示文稿",
            "",
            "设计理念",
            "  数学严谨性 + 自动化生成 + 教学实用性",
        ], band_color="coral")

        # ---- Slide 4: Architecture ----
        self._slide_content("2. 系统架构", bullets_left=[
            "模块划分",
            "",
            "1. ProblemGenerator（题目生成器）",
            "   椭圆题目生成 (3个难度)",
            "   双曲线题目生成 (3个难度)",
            "   抛物线题目生成 (3个难度)",
            "   极坐标题目生成 (3个难度)",
            "",
            "2. DiagramRenderer（配图渲染器）",
            "   坐标系建立（网格、刻度、箭头）",
            "   圆锥曲线绘制（参数方程，1000采样点）",
            "   关键点/线标注 + LaTeX标签",
        ], bullets_right=[
            "数据结构",
            "",
            "ConicParams（圆锥曲线参数）",
            "  center: 中心点坐标",
            "  a, b: 半轴长度",
            "  c: 半焦距",
            "  e: 离心率",
            "",
            "Problem（题目对象）",
            "  problem_latex: LaTeX题干",
            "  solution_latex: LaTeX解答",
            "  points: 关键点列表",
            "  lines: 关键直线列表",
            "  conic_type: 曲线类型",
        ], band_color="dark")

        # ---- Slide 5: Algorithm — Problem Generation ----
        self._slide_content("3. 核心算法：题目生成引擎", bullets_left=[
            "设计原则",
            "  参数随机化：在合理范围内随机生成几何参数",
            "  数学可解性：所有生成的题目都有明确的解析解",
            "  渐进式难度：从基础概念到综合应用，难度递增",
            "",
            "核心算法流程",
            "  Step 1: 参数采样 — 随机生成 a, b 等基本参数",
            "  Step 2: 参数推导 — 计算 c, e 等派生参数",
            "  Step 3: 几何对象构造 — 生成焦点、顶点、渐近线等",
            "  Step 4: 交点计算 — 联立方程求解弦与曲线的交点",
            "  Step 5: 题干生成 — 将参数代入LaTeX模板",
            "  Step 6: 解答推导 — 自动计算并生成完整解答过程",
        ], bullets_right=[
            "椭圆生成示例",
            "",
            "1. 采样 a ∈ {2,3,4,5,6}, b < a",
            "2. c = √(a² − b²),  e = c/a",
            "3. 焦点 F₁(−c,0), F₂(c,0)",
            "4. 焦点弦: 联立 y=k(x+c) 与 x²/a²+y²/b²=1",
            "5. 韦达定理得 x₁+x₂, x₁x₂",
            "6. 弦长 |PQ| = √(1+k²)·|x₁−x₂|",
            "7. 面积 S = ½·|PQ|·d(F₁, 直线)",
            "",
            "所有中间量精确计算，无近似",
        ], band_color="coral")

        # ---- Slide 6: Algorithm — Diagram Rendering ----
        self._slide_content("4. 核心算法：精确配图渲染", bullets_left=[
            "渲染流程",
            "  Step 1: 坐标系建立 — 轴范围、网格、刻度、箭头",
            "  Step 2: 曲线绘制 — 参数方程精确绘制",
            "  Step 3: 直线绘制 — ax+by+c=0",
            "  Step 4: 点标注 — 精确位置绘制关键点",
            "  Step 5: 标签渲染 — LaTeX格式点标签",
            "  Step 6: 图形导出 — 高分辨率 PNG",
            "",
            "精确性保证",
            "  所有曲线使用参数方程绘制，1000个采样点",
            "  关键点坐标由精确计算得到，非估算",
            "  坐标轴等比例显示，保证几何关系视觉正确",
        ], bullets_right=[
            "配色方案",
            "  蓝色 #2E86AB — 圆锥曲线",
            "  红色 #E8451E — 直线、弦",
            "  紫色 #9B59B6 — 渐近线（虚线）",
            "  绿色 #27AE60 — 准线（点划线）",
            "  金色 #FFD700 — 关键点填充",
            "",
            "渲染特性",
            "  10:1 等比例坐标轴",
            "  箭头标注坐标轴方向",
            "  半透明曲线填充区域",
            "  焦点特殊颜色标记（金色填充）",
        ], band_color="dark")

        # ---- Slides 7-10: Demo images ----
        ellipse_imgs = [
            "/root/analytic_geometry_generator/output/椭圆_difficulty1.png",
            "/root/analytic_geometry_generator/output/椭圆_difficulty2.png",
            "/root/analytic_geometry_generator/output/椭圆_difficulty3.png",
        ]
        self._slide_images("5. 演示：椭圆题目",
                           ellipse_imgs,
                           ["基础题：椭圆方程与焦点",
                            "进阶题：焦点弦长与面积",
                            "竞赛题：焦点三角形"])

        hyperbola_imgs = [
            "/root/analytic_geometry_generator/output/双曲线_difficulty1.png",
            "/root/analytic_geometry_generator/output/双曲线_difficulty2.png",
            "/root/analytic_geometry_generator/output/双曲线_difficulty3.png",
        ]
        self._slide_images("6. 演示：双曲线题目",
                           hyperbola_imgs,
                           ["基础题：双曲线方程与渐近线",
                            "进阶题：焦点弦",
                            "竞赛题：焦点三角形"])

        parabola_imgs = [
            "/root/analytic_geometry_generator/output/抛物线_difficulty1.png",
            "/root/analytic_geometry_generator/output/抛物线_difficulty2.png",
            "/root/analytic_geometry_generator/output/抛物线_difficulty3.png",
        ]
        self._slide_images("7. 演示：抛物线题目",
                           parabola_imgs,
                           ["基础题：抛物线方程与焦点",
                            "进阶题：焦点弦",
                            "竞赛题：焦点弦性质证明"])

        polar_imgs = [
            "/root/analytic_geometry_generator/output/极坐标_difficulty1.png",
            "/root/analytic_geometry_generator/output/极坐标_difficulty2.png",
            "/root/analytic_geometry_generator/output/极坐标_difficulty3.png",
        ]
        self._slide_images("8. 演示：极坐标题目",
                           polar_imgs,
                           ["基础题：极坐标与直角坐标互化",
                            "进阶题：直线与圆的位置关系",
                            "竞赛题：圆锥曲线极坐标方程"])

        # ---- Slide 11: Highlights ----
        self._slide_content("9. 技术亮点与扩展方向", bullets_left=[
            "✓ 数学严谨性 — 所有题目均可解，配图与参数严格一致",
            "✓ 自动化生成 — 无需人工干预，一键生成完整题目",
            "✓ LaTeX输出 — 标准数学排版格式，便于教学使用",
            "✓ 高质量配图 — Matplotlib渲染，支持矢量图导出",
            "✓ 模块化设计 — 易于扩展新的题目类型",
            "",
            "架构优势",
            "  ProblemGenerator 与 DiagramRenderer 完全解耦",
            "  同一 Problem 对象可渲染为 PNG / TikZ / SVG",
            "  随机种子保证可复现",
        ], bullets_right=[
            "扩展方向",
            "  → 更多知识点：圆、参数方程、直线方程",
            "  → 自定义难度：更细粒度的难度控制",
            "  → 题目数据库：生成的题目持久化存储",
            "  → Web界面：在线题目生成平台",
            "  → TikZ输出：LaTeX TikZ绘图代码",
            "  → 智能组卷：知识点覆盖率自动组卷",
            "",
            "技术栈",
            "  Python 3.12 + NumPy + Matplotlib",
            "  python-pptx (演示文稿生成)",
            "  LaTeX (数学排版)",
        ], band_color="coral")

        # ---- Slide 12: CTA ----
        self._slide_cta_coral()

        # Save
        self.prs.save(output_path)
        return output_path


def main():
    """Create the PPT"""
    creator = GeometryPPTCreator()
    path = creator.create_presentation()
    print(f"PPT已生成: {path}")


if __name__ == "__main__":
    main()
