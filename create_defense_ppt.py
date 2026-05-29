"""
解析几何题目生成系统 — 答辩演示文稿
Oral Defense Presentation

基于 Anthropic DESIGN.md 设计规范
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os


class DefensePPTCreator:
    """答辩PPT创建器 — Claude Design System"""

    # ==================== DESIGN TOKENS ====================
    CANVAS = RGBColor(0xFA, 0xF9, 0xF5)
    SURFACE_CARD = RGBColor(0xEF, 0xE9, 0xDE)
    SURFACE_DARK = RGBColor(0x18, 0x17, 0x15)
    SURFACE_DARK_ELEVATED = RGBColor(0x25, 0x23, 0x20)
    PRIMARY = RGBColor(0xCC, 0x78, 0x5C)
    PRIMARY_ACTIVE = RGBColor(0xA9, 0x58, 0x3E)
    ON_PRIMARY = RGBColor(0xFF, 0xFF, 0xFF)
    INK = RGBColor(0x14, 0x14, 0x13)
    BODY = RGBColor(0x3D, 0x3D, 0x3A)
    BODY_STRONG = RGBColor(0x25, 0x25, 0x23)
    MUTED = RGBColor(0x6C, 0x6A, 0x64)
    MUTED_SOFT = RGBColor(0x8E, 0x8B, 0x82)
    HAIRLINE = RGBColor(0xE6, 0xDF, 0xD8)
    ON_DARK = RGBColor(0xFA, 0xF9, 0xF5)
    ON_DARK_SOFT = RGBColor(0xA0, 0x9D, 0x96)
    ACCENT_TEAL = RGBColor(0x5D, 0xB8, 0xA6)
    ACCENT_AMBER = RGBColor(0xE8, 0xA5, 0x5A)
    SUCCESS = RGBColor(0x5D, 0xB8, 0x72)

    DISPLAY_FONT = "Georgia"
    BODY_FONT = "Calibri"

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    # ==================== HELPERS ====================

    def _set_bg(self, slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _text(self, slide, left, top, w, h, text,
              font=None, size=Pt(16), color=None,
              bold=False, italic=False, align=PP_ALIGN.LEFT,
              line_spacing=1.5, space_after=Pt(6)):
        if font is None: font = self.BODY_FONT
        if color is None: color = self.BODY
        tb = slide.shapes.add_textbox(left, top, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = space_after
        run = p.runs[0]
        run.font.name = font
        run.font.size = size
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        return tb, tf

    def _heading(self, slide, left, top, w, h, text, size="lg", color=None):
        smap = {"xl": (Pt(64), 1.05), "lg": (Pt(48), 1.1),
                "md": (Pt(36), 1.15), "sm": (Pt(28), 1.2)}
        fs, ls = smap.get(size, (Pt(48), 1.1))
        if color is None: color = self.INK
        return self._text(slide, left, top, w, h, text,
                          font=self.DISPLAY_FONT, size=fs, color=color,
                          bold=False, line_spacing=ls, space_after=Pt(4))

    def _subheading(self, slide, left, top, w, h, text, size="md", color=None):
        smap = {"lg": Pt(22), "md": Pt(18), "sm": Pt(16)}
        fs = smap.get(size, Pt(18))
        if color is None: color = self.BODY_STRONG
        return self._text(slide, left, top, w, h, text,
                          font=self.BODY_FONT, size=fs, color=color,
                          bold=True, line_spacing=1.4, space_after=Pt(4))

    def _bullets(self, slide, left, top, w, h, items,
                 font=None, size=Pt(16), color=None,
                 spacing=1.5, bullet="•", after=Pt(5)):
        if font is None: font = self.BODY_FONT
        if color is None: color = self.BODY
        tb = slide.shapes.add_textbox(left, top, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"{bullet}  {item}" if bullet else item
            p.line_spacing = spacing
            p.space_after = after
            run = p.runs[0]
            run.font.name = font
            run.font.size = size
            run.font.color.rgb = color
        return tb

    def _band(self, slide, color, top=0, h=Inches(1.2)):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0), top, Inches(13.333), h)
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        return s

    def _card(self, slide, left, top, w, h, bg=None, border=True):
        if bg is None: bg = self.SURFACE_CARD
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = bg
        if border:
            s.line.color.rgb = self.HAIRLINE
            s.line.width = Pt(1)
        else:
            s.line.fill.background()
        return s

    def _dark_card(self, slide, left, top, w, h):
        return self._card(slide, left, top, w, h, bg=self.SURFACE_DARK, border=False)

    # ==================== SLIDE BUILDERS ====================

    def _slide_cover(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_bg(slide, self.CANVAS)

        # Coral accent stripe
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(1), Inches(1.8),
                                        Inches(0.08), Inches(3.5))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = self.PRIMARY
        stripe.line.fill.background()

        # Title
        self._heading(slide, Inches(1.4), Inches(1.8), Inches(10), Inches(1.5),
                      "解析几何题目生成系统", size="xl")

        # Subtitle
        self._text(slide, Inches(1.4), Inches(3.5), Inches(9), Inches(0.7),
                   "技术原理与实现",
                   font=self.DISPLAY_FONT, size=Pt(28), color=self.MUTED,
                   line_spacing=1.2)

        # Meta info card
        card = self._card(slide, Inches(1.4), Inches(4.6), Inches(5.5), Inches(1.8))
        self._text(slide, Inches(1.7), Inches(4.8), Inches(5), Inches(0.4),
                   "项目类型", size=Pt(12), color=self.MUTED, bold=True)
        self._text(slide, Inches(1.7), Inches(5.1), Inches(5), Inches(0.4),
                   "解析几何自动化出题与配图系统", size=Pt(16), color=self.BODY_STRONG, bold=True)
        self._text(slide, Inches(1.7), Inches(5.5), Inches(5), Inches(0.4),
                   "技术栈：Python 3 + NumPy + Matplotlib + python-pptx",
                   size=Pt(13), color=self.MUTED)

        # Coral badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(1.4), Inches(6.6),
                                       Inches(2.8), Inches(0.5))
        badge.fill.solid()
        badge.fill.fore_color.rgb = self.PRIMARY
        badge.line.fill.background()
        self._text(slide, Inches(1.4), Inches(6.62), Inches(2.8), Inches(0.45),
                   "DESIGN SYSTEM", size=Pt(13), color=self.ON_PRIMARY,
                   bold=True, align=PP_ALIGN.CENTER)

    def _slide_outline(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_bg(slide, self.CANVAS)

        self._heading(slide, Inches(1), Inches(0.6), Inches(11), Inches(1),
                      "汇报提纲", size="lg")

        div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(1), Inches(1.5), Inches(2), Inches(0.04))
        div.fill.solid()
        div.fill.fore_color.rgb = self.PRIMARY
        div.line.fill.background()

        sections = [
            ("01", "项目背景与需求分析"),
            ("02", "系统总体架构设计"),
            ("03", "核心算法：数学原理"),
            ("04", "核心算法：题目生成引擎"),
            ("05", "核心算法：精确配图渲染"),
            ("06", "系统演示"),
            ("07", "技术亮点与创新点"),
            ("08", "测试与验证"),
            ("09", "总结与展望"),
        ]

        for i, (num, title) in enumerate(sections):
            y = Inches(1.9) + Inches(i * 0.58)
            self._text(slide, Inches(1), y, Inches(0.8), Inches(0.45),
                       num, font=self.DISPLAY_FONT, size=Pt(22),
                       color=self.PRIMARY, bold=True, line_spacing=1.0)
            self._text(slide, Inches(1.9), y, Inches(9), Inches(0.45),
                       title, size=Pt(18), color=self.BODY_STRONG, line_spacing=1.0)

    def _slide_content(self, title, left_items=None, right_items=None,
                       image=None, band="coral"):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_bg(slide, self.CANVAS)

        bc = self.SURFACE_DARK if band == "dark" else self.PRIMARY
        self._band(slide, bc)
        tc = self.ON_DARK if band == "dark" else self.ON_PRIMARY
        self._text(slide, Inches(0.8), Inches(0.2), Inches(11.5), Inches(0.8),
                   title, font=self.DISPLAY_FONT, size=Pt(36),
                   color=tc, bold=False, line_spacing=1.15, space_after=Pt(0))

        if image and os.path.exists(image):
            self._dark_card(slide, Inches(7.2), Inches(1.6), Inches(5.5), Inches(5.3))
            slide.shapes.add_picture(image, Inches(7.4), Inches(1.8),
                                     Inches(5.1), Inches(4.9))
            if left_items:
                self._card(slide, Inches(0.6), Inches(1.6), Inches(6.2), Inches(5.3))
                self._bullets(slide, Inches(0.9), Inches(1.9), Inches(5.6), Inches(4.8),
                              left_items, size=Pt(15))
        elif left_items and right_items:
            self._card(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.3))
            self._bullets(slide, Inches(0.9), Inches(1.9), Inches(5.2), Inches(4.8),
                          left_items, size=Pt(15))
            self._card(slide, Inches(6.9), Inches(1.6), Inches(5.8), Inches(5.3))
            self._bullets(slide, Inches(7.2), Inches(1.9), Inches(5.2), Inches(4.8),
                          right_items, size=Pt(15))
        elif left_items:
            self._card(slide, Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.3))
            self._bullets(slide, Inches(0.9), Inches(1.9), Inches(11.5), Inches(4.8),
                          left_items, size=Pt(15))

    def _slide_images(self, title, images, captions=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_bg(slide, self.CANVAS)
        self._band(slide, self.SURFACE_DARK)
        self._text(slide, Inches(0.8), Inches(0.2), Inches(11.5), Inches(0.8),
                   title, font=self.DISPLAY_FONT, size=Pt(36),
                   color=self.ON_DARK, bold=False, line_spacing=1.15)

        n = len(images)
        if n == 1:
            pos = [(Inches(2.5), Inches(1.6), Inches(8.3), Inches(5.3))]
        elif n == 2:
            pos = [(Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.3)),
                   (Inches(6.9), Inches(1.6), Inches(5.8), Inches(5.3))]
        else:
            pos = [(Inches(0.6), Inches(1.6), Inches(3.8), Inches(5.3)),
                   (Inches(4.7), Inches(1.6), Inches(3.8), Inches(5.3)),
                   (Inches(8.9), Inches(1.6), Inches(3.8), Inches(5.3))]

        for i, img in enumerate(images[:3]):
            if i < len(pos) and os.path.exists(img):
                l, t, w, h = pos[i]
                self._dark_card(slide, l, t, w, h)
                slide.shapes.add_picture(img, l + Inches(0.15), t + Inches(0.15),
                                         w - Inches(0.3), h - Inches(0.8))
                if captions and i < len(captions):
                    self._text(slide, l, t + h - Inches(0.55), w, Inches(0.45),
                               captions[i], size=Pt(12), color=self.MUTED_SOFT,
                               align=PP_ALIGN.CENTER, line_spacing=1.0)

    def _slide_formula(self, title, formulas, explanation=None):
        """公式展示页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_bg(slide, self.CANVAS)
        self._band(slide, self.PRIMARY)
        self._text(slide, Inches(0.8), Inches(0.2), Inches(11.5), Inches(0.8),
                   title, font=self.DISPLAY_FONT, size=Pt(36),
                   color=self.ON_PRIMARY, bold=False, line_spacing=1.15)

        # Dark card for formulas
        self._dark_card(slide, Inches(0.6), Inches(1.6), Inches(12.1), Inches(3.2))
        for i, formula in enumerate(formulas):
            self._text(slide, Inches(1.2), Inches(1.9 + i * 0.6), Inches(11), Inches(0.5),
                       formula, font="Consolas", size=Pt(18),
                       color=self.ACCENT_TEAL, line_spacing=1.3)

        if explanation:
            self._card(slide, Inches(0.6), Inches(5.1), Inches(12.1), Inches(1.8))
            self._bullets(slide, Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.5),
                          explanation, size=Pt(15))

    def _slide_thanks(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_bg(slide, self.CANVAS)

        # Coral card
        self._card(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.1),
                   bg=self.PRIMARY, border=False)

        self._text(slide, Inches(1.6), Inches(2.0), Inches(10), Inches(1.5),
                   "谢谢各位老师",
                   font=self.DISPLAY_FONT, size=Pt(64),
                   color=self.ON_PRIMARY, bold=False, line_spacing=1.05,
                   align=PP_ALIGN.CENTER)

        self._text(slide, Inches(1.6), Inches(3.8), Inches(10), Inches(1),
                   "请各位老师批评指正",
                   size=Pt(24), color=RGBColor(0xFF, 0xE0, 0xD0),
                   align=PP_ALIGN.CENTER, line_spacing=1.3)

        btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(5.0), Inches(5.0), Inches(3.3), Inches(0.5))
        btn.fill.solid()
        btn.fill.fore_color.rgb = self.CANVAS
        btn.line.fill.background()
        self._text(slide, Inches(5.0), Inches(5.02), Inches(3.3), Inches(0.45),
                   "Q & A", font=self.DISPLAY_FONT, size=Pt(18),
                   color=self.PRIMARY, bold=True, align=PP_ALIGN.CENTER)

    # ==================== BUILD ====================

    def create_presentation(self, output_path=None):
        if output_path is None:
            output_path = "/root/analytic_geometry_generator/output/答辩PPT.pptx"
        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        # 1. Cover
        self._slide_cover()

        # 2. Outline
        self._slide_outline()

        # 3. Background
        self._slide_content("一、项目背景与需求分析", left_items=[
            "传统出题方式的痛点",
            "  手工出题效率低，需大量计算验证",
            "  手工绘图不精确，与数学参数不一致",
            "  题型单一，受教师经验限制",
            "  难以快速生成大量变式题目",
            "",
            "项目目标",
            "  根据知识点与难度自动生成解析几何题目",
            "  所有题目数学可解，解答完整可靠",
            "  精确配图，图形与参数严格一致",
            "  LaTeX 标准排版，便于教学使用",
        ], right_items=[
            "适用场景",
            "  教师日常出题与组卷",
            "  学生自主练习",
            "  教学资源开发",
            "",
            "支持的知识点",
            "  椭圆 (Ellipse)",
            "  双曲线 (Hyperbola)",
            "  抛物线 (Parabola)",
            "  极坐标 (Polar Coordinates)",
            "",
            "难度等级",
            "  ★☆☆ 基础 → ★★☆ 进阶 → ★★★ 竞赛",
        ], band="coral")

        # 4. Architecture
        self._slide_content("二、系统总体架构设计", left_items=[
            "三层架构",
            "",
            "1. 主程序层 (main.py)",
            "   协调各模块，提供统一入口",
            "",
            "2. 核心引擎层",
            "   ProblemGenerator — 题目生成",
            "   DiagramRenderer — 配图渲染",
            "   PPT Creator — 演示文稿",
            "",
            "3. 输出层",
            "   LaTeX 题干 + 解答",
            "   PNG 精确配图",
            "   PPTX 演示文稿",
        ], right_items=[
            "数据结构设计",
            "",
            "ConicParams — 圆锥曲线参数",
            "  center, a, b, c, e",
            "",
            "Point — 关键点",
            "  x, y, label (LaTeX)",
            "",
            "Line — 直线",
            "  ax + by + c = 0",
            "",
            "Problem — 题目对象",
            "  串联所有数据",
            "  包含 LaTeX 文本",
            "  包含几何参数",
        ], band="dark")

        # 5. Math principles
        self._slide_formula("三、核心算法：数学原理", [
            "椭圆:  x²/a² + y²/b² = 1,    c = √(a² − b²),  e = c/a < 1",
            "双曲线: x²/a² − y²/b² = 1,    c = √(a² + b²),  e = c/a > 1",
            "抛物线: y² = 2px,              焦点 F(p/2, 0),  准线 x = −p/2",
            "极坐标: ρ = ep/(1 − e·cosθ),  圆: ρ = 2r·cosθ",
        ], explanation=[
            "椭圆焦点弦：联立 y = k(x+c) 与椭圆方程，韦达定理求交点",
            "双曲线渐近线：y = ±(b/a)x，由 x²/a² − y²/b² = 0 推导",
            "抛物线焦点弦：|PQ| = 2p/sin²θ，1/|PF| + 1/|QF| = 2/p (定值)",
            "极坐标互化：ρ² = x² + y²,  x = ρcosθ,  y = ρsinθ",
        ])

        # 6. Problem generation
        self._slide_content("四、核心算法：题目生成引擎", left_items=[
            "设计原则",
            "  参数随机化 — 合理范围内随机生成几何参数",
            "  数学可解性 — 所有题目有明确解析解",
            "  渐进式难度 — 从基础到竞赛，逐级递增",
            "",
            "生成流程",
            "  Step 1: 参数采样 (a, b, c, e)",
            "  Step 2: 参数推导 (派生量计算)",
            "  Step 3: 几何对象构造 (焦点、顶点等)",
            "  Step 4: 交点计算 (联立方程)",
            "  Step 5: 题干生成 (LaTeX 模板)",
            "  Step 6: 解答推导 (完整过程)",
        ], right_items=[
            "参数采样约束",
            "  椭圆: a > b > 0",
            "  双曲线: a > 0, b > 0",
            "  抛物线: p > 0",
            "",
            "可解性保证",
            "  判别式 Δ ≥ 0 → 实交点存在",
            "  韦达定理 → 弦长/面积精确",
            "  浮点运算保留 4 位有效数字",
            "",
            "难度梯度",
            "  Level 1: 标准方程、基本参数",
            "  Level 2: 弦长、面积、位置关系",
            "  Level 3: 综合证明、最值问题",
        ], band="coral")

        # 7. Rendering
        self._slide_content("五、核心算法：精确配图渲染", left_items=[
            "渲染流程",
            "  Step 1: 坐标系建立 (轴、网格、箭头)",
            "  Step 2: 曲线绘制 (参数方程, 1000点)",
            "  Step 3: 直线绘制 (ax+by+c=0)",
            "  Step 4: 点标注 (精确位置)",
            "  Step 5: 标签渲染 (LaTeX 标签)",
            "  Step 6: 图形导出 (高分辨率 PNG)",
            "",
            "精确性保证",
            "  参数方程绘制，非近似",
            "  1000 采样点保证光滑",
            "  等比例坐标轴",
        ], right_items=[
            "颜色编码系统",
            "  蓝色 — 圆锥曲线",
            "  红色 — 直线、弦",
            "  紫色虚线 — 渐近线",
            "  绿色点划线 — 准线",
            "  金色 — 关键点填充",
            "",
            "z-order 层次",
            "  Layer 1: 坐标轴、网格",
            "  Layer 2: 曲线填充",
            "  Layer 3: 曲线轮廓",
            "  Layer 4: 直线",
            "  Layer 5: 点标注 + 标签",
        ], band="dark")

        # 8. Demo slides
        self._slide_images("六、系统演示 — 椭圆", [
            "/root/analytic_geometry_generator/output/椭圆_difficulty1.png",
            "/root/analytic_geometry_generator/output/椭圆_difficulty2.png",
            "/root/analytic_geometry_generator/output/椭圆_difficulty3.png",
        ], ["基础：椭圆方程与焦点", "进阶：焦点弦长与面积", "竞赛：焦点三角形"])

        self._slide_images("六、系统演示 — 双曲线", [
            "/root/analytic_geometry_generator/output/双曲线_difficulty1.png",
            "/root/analytic_geometry_generator/output/双曲线_difficulty2.png",
            "/root/analytic_geometry_generator/output/双曲线_difficulty3.png",
        ], ["基础：双曲线方程与渐近线", "进阶：焦点弦", "竞赛：焦点三角形"])

        self._slide_images("六、系统演示 — 抛物线与极坐标", [
            "/root/analytic_geometry_generator/output/抛物线_difficulty2.png",
            "/root/analytic_geometry_generator/output/极坐标_difficulty2.png",
        ], ["抛物线进阶：焦点弦", "极坐标进阶：直线与圆"])

        # 9. Innovation
        self._slide_content("七、技术亮点与创新点", left_items=[
            "创新点 1: 参数化题目生成",
            "  通过随机采样 + 约束条件",
            "  自动生成多样化可解题目",
            "",
            "创新点 2: 数学可解性保证",
            "  所有中间量精确计算",
            "  解答过程完整可靠",
            "",
            "创新点 3: 图文严格一致",
            "  配图参数与题干数学参数",
            "  由同一 Problem 对象驱动",
        ], right_items=[
            "创新点 4: 模块化解耦设计",
            "  ProblemGenerator 与 DiagramRenderer",
            "  完全解耦，可独立扩展",
            "",
            "创新点 5: 多格式输出",
            "  LaTeX 文本 + PNG 图片 + PPTX 演示",
            "  同一数据源，多种表现形式",
            "",
            "技术亮点",
            "  Matplotlib 参数方程精确绘制",
            "  LaTeX 标准数学排版",
            "  随机种子保证可复现",
        ], band="coral")

        # 10. Testing
        self._slide_content("八、测试与验证", left_items=[
            "参数一致性验证",
            "  椭圆参数: a > b > 0  ✓",
            "  焦距关系: c = √(a² − b²)  ✓",
            "  离心率: e = c/a ∈ (0,1)  ✓",
            "  交点判别式: Δ ≥ 0  ✓",
            "  韦达定理: x₁+x₂, x₁x₂ 正确  ✓",
            "",
            "配图精度验证",
            "  1000 采样点，视觉光滑  ✓",
            "  关键点坐标精确计算  ✓",
            "  坐标轴等比例显示  ✓",
        ], right_items=[
            "LaTeX 格式验证",
            "  数学公式正确嵌套  ✓",
            "  花括号转义正确  ✓",
            "  符号显示正确  ✓",
            "  答案格式完整  ✓",
            "",
            "测试覆盖",
            "  4 个知识点 × 3 个难度 = 12 道题",
            "  每道题含完整题干 + 解答 + 配图",
            "  所有配图通过视觉检查",
            "  所有数学公式通过 LaTeX 渲染验证",
        ], band="dark")

        # 11. Summary
        self._slide_content("九、总结与展望", left_items=[
            "已完成工作",
            "  ✓ 题目生成引擎 (12道题, 4知识点×3难度)",
            "  ✓ 精确配图渲染器 (12张高精度配图)",
            "  ✓ LaTeX 输出 (题干+完整解答)",
            "  ✓ PPT 演示文稿生成",
            "  ✓ 答辩演示文稿",
            "",
            "技术栈",
            "  Python 3.12 + NumPy + Matplotlib",
            "  python-pptx + LaTeX",
        ], right_items=[
            "未来展望",
            "  → 更多知识点: 圆、参数方程、直线方程",
            "  → TikZ 输出: LaTeX TikZ 绘图代码",
            "  → Web 界面: 在线题目生成平台",
            "  → 题目数据库: 生成题目持久化存储",
            "  → 智能组卷: 知识点覆盖率自动组卷",
            "",
            "核心价值",
            "  数学严谨性 + 自动化 + 教学实用性",
        ], band="coral")

        # 12. Thanks
        self._slide_thanks()

        self.prs.save(output_path)
        return output_path


def main():
    creator = DefensePPTCreator()
    path = creator.create_presentation()
    print(f"答辩PPT已生成: {path}")


if __name__ == "__main__":
    main()
